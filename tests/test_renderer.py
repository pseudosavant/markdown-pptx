from __future__ import annotations

import base64
import shutil
import struct
import xml.etree.ElementTree as ET
import zipfile
import zlib
from copy import deepcopy
from pathlib import Path

import httpx
import pytest
from pptx import Presentation

from markdown_slides.assets import default_template_path
from markdown_slides.errors import AssetError, RenderError, TemplateError
from markdown_slides.parser import parse_deck
from markdown_slides.renderer import Downloader, list_layout_details, list_master_details, render_pptx

PNG_BYTES = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAIAAACQd1PeAAAADElEQVR4nGP4//8/AAX+Av4N70a4AAAAAElFTkSuQmCC"
)
NS = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
TABLE_STYLE_MEDIUM_1_ACCENT_1 = "{B301B821-A1FF-4177-AEE7-76D212191A09}"


def _png_header(*, width: int, height: int) -> bytes:
    def chunk(chunk_type: bytes, data: bytes) -> bytes:
        crc = zlib.crc32(chunk_type + data) & 0xFFFFFFFF
        return struct.pack(">I", len(data)) + chunk_type + data + struct.pack(">I", crc)

    header = struct.pack(">IIBBBBB", width, height, 8, 2, 0, 0, 0)
    return b"\x89PNG\r\n\x1a\n" + chunk(b"IHDR", header) + chunk(b"IDAT", zlib.compress(b"")) + chunk(b"IEND", b"")


class FakeDownloader(Downloader):
    def fetch(self, url: str) -> bytes:
        return PNG_BYTES


def _make_two_master_template(path: Path) -> None:
    content_type_ns = "http://schemas.openxmlformats.org/package/2006/content-types"
    relationship_ns = "http://schemas.openxmlformats.org/package/2006/relationships"
    presentation_ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
    document_rel_ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    source_parts: dict[str, bytes] = {}
    copied_parts = {
        "ppt/slideMasters/slideMaster1.xml",
        "ppt/slideMasters/_rels/slideMaster1.xml.rels",
        "ppt/slideLayouts/slideLayout2.xml",
        "ppt/slideLayouts/_rels/slideLayout2.xml.rels",
        "ppt/theme/theme1.xml",
    }
    with (
        zipfile.ZipFile(default_template_path(), "r") as source,
        zipfile.ZipFile(path, "w", compression=zipfile.ZIP_DEFLATED) as target,
    ):
        for info in source.infolist():
            data = source.read(info.filename)
            if info.filename in copied_parts:
                source_parts[info.filename] = data
            if info.filename == "[Content_Types].xml":
                xml = ET.fromstring(data)
                for part_name, content_type in (
                    (
                        "/ppt/slideMasters/slideMaster2.xml",
                        "application/vnd.openxmlformats-officedocument.presentationml.slideMaster+xml",
                    ),
                    (
                        "/ppt/slideLayouts/slideLayout12.xml",
                        "application/vnd.openxmlformats-officedocument.presentationml.slideLayout+xml",
                    ),
                    ("/ppt/theme/theme7.xml", "application/vnd.openxmlformats-officedocument.theme+xml"),
                ):
                    ET.SubElement(
                        xml,
                        f"{{{content_type_ns}}}Override",
                        {"PartName": part_name, "ContentType": content_type},
                    )
                data = ET.tostring(xml, encoding="utf-8", xml_declaration=True)
            elif info.filename == "ppt/_rels/presentation.xml.rels":
                xml = ET.fromstring(data)
                ET.SubElement(
                    xml,
                    f"{{{relationship_ns}}}Relationship",
                    {
                        "Id": "rId999",
                        "Type": "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideMaster",
                        "Target": "slideMasters/slideMaster2.xml",
                    },
                )
                data = ET.tostring(xml, encoding="utf-8", xml_declaration=True)
            elif info.filename == "ppt/presentation.xml":
                xml = ET.fromstring(data)
                master_ids = xml.find(f"{{{presentation_ns}}}sldMasterIdLst")
                ET.SubElement(
                    master_ids,
                    f"{{{presentation_ns}}}sldMasterId",
                    {"id": "2147483650", f"{{{document_rel_ns}}}id": "rId999"},
                )
                data = ET.tostring(xml, encoding="utf-8", xml_declaration=True)
            target.writestr(info, data)

        assert copied_parts == source_parts.keys()
        master = ET.fromstring(source_parts["ppt/slideMasters/slideMaster1.xml"])
        master.find(f"{{{presentation_ns}}}cSld").set("name", "Executive")
        target.writestr(
            "ppt/slideMasters/slideMaster2.xml",
            ET.tostring(master, encoding="utf-8", xml_declaration=True),
        )

        master_rels = ET.fromstring(source_parts["ppt/slideMasters/_rels/slideMaster1.xml.rels"])
        for relationship in master_rels.findall(f"{{{relationship_ns}}}Relationship"):
            if relationship.attrib["Type"].endswith("/theme"):
                relationship.set("Target", "../theme/theme7.xml")
            elif relationship.attrib["Target"].endswith("/slideLayout2.xml"):
                relationship.set("Target", "../slideLayouts/slideLayout12.xml")
        target.writestr(
            "ppt/slideMasters/_rels/slideMaster2.xml.rels",
            ET.tostring(master_rels, encoding="utf-8", xml_declaration=True),
        )

        target.writestr("ppt/slideLayouts/slideLayout12.xml", source_parts["ppt/slideLayouts/slideLayout2.xml"])
        layout_rels = ET.fromstring(source_parts["ppt/slideLayouts/_rels/slideLayout2.xml.rels"])
        for relationship in layout_rels.findall(f"{{{relationship_ns}}}Relationship"):
            if relationship.attrib["Type"].endswith("/slideMaster"):
                relationship.set("Target", "../slideMasters/slideMaster2.xml")
        target.writestr(
            "ppt/slideLayouts/_rels/slideLayout12.xml.rels",
            ET.tostring(layout_rels, encoding="utf-8", xml_declaration=True),
        )

        theme = ET.fromstring(source_parts["ppt/theme/theme1.xml"])
        theme.set("name", "Executive Theme")
        target.writestr("ppt/theme/theme7.xml", ET.tostring(theme, encoding="utf-8", xml_declaration=True))


def test_render_preserves_text_placeholders_and_notes(tmp_path: Path) -> None:
    deck = parse_deck(
        """---
fonts:
  body: Aptos
  headings: Aptos Display
title_color: "#112233"
body_color: "#445566"
color_scheme:
  preset: Office
---

# Title slide title
---
layout: Title Slide
body_color: "#778899"
notes: |
  Note line.
---

Title slide subtitle

# Title and content title

Title and content text

# Section header title
---
layout: Section Header
---

Section header subtitle
""",
        input_path=tmp_path / "deck.md",
        source_name=str(tmp_path / "deck.md"),
    )
    output = tmp_path / "deck.pptx"

    render_pptx(deck, output_path=output, template_path=None, force=False, base_dir=tmp_path)

    prs = Presentation(str(output))
    assert len(prs.slides) == 3
    assert prs.slides[0].shapes.title.text == "Title slide title"
    assert prs.slides[0].placeholders[1].text == "Title slide subtitle"
    assert prs.slides[1].shapes.title.text == "Title and content title"
    assert "Title and content text" in prs.slides[1].placeholders[1].text
    assert prs.slides[2].shapes.title.text == "Section header title"
    assert prs.slides[2].placeholders[1].text == "Section header subtitle"
    assert prs.slides[0].notes_slide.notes_text_frame.text == "Note line."
    with zipfile.ZipFile(output) as zf:
        slide1_xml = zf.read("ppt/slides/slide1.xml").decode("utf-8")
        slide2_xml = zf.read("ppt/slides/slide2.xml").decode("utf-8")
    assert 'val="112233"' in slide1_xml
    assert 'val="778899"' in slide1_xml
    assert 'val="112233"' in slide2_xml
    assert 'val="445566"' in slide2_xml


def test_render_sets_theme_and_aspect_ratio(tmp_path: Path) -> None:
    deck = parse_deck(
        """---
aspect_ratio: "4:3"
fonts:
  body: Aptos
  headings: Aptos Display
color_scheme:
  preset: Blue Warm
---

# Slide

Body
""",
        input_path=tmp_path / "deck.md",
        source_name=str(tmp_path / "deck.md"),
    )
    output = tmp_path / "deck.pptx"

    render_pptx(deck, output_path=output, template_path=None, force=False, base_dir=tmp_path)

    prs = Presentation(str(output))
    assert prs.slide_width == 9144000
    assert prs.slide_height == 6858000
    with zipfile.ZipFile(output) as zf:
        theme = ET.fromstring(zf.read("ppt/theme/theme1.xml"))
    clr = theme.find(".//a:clrScheme", NS)
    assert clr is not None
    assert clr.attrib["name"] == "Blue Warm"
    accent1 = clr.find("a:accent1", NS)[0]
    assert accent1.attrib["val"] == "4A66AC"


def test_render_preserves_custom_template_theme_when_markdown_has_no_overrides(tmp_path: Path) -> None:
    template = tmp_path / "template.pptx"
    shutil.copyfile(default_template_path(), template)
    with (
        zipfile.ZipFile(template, "r") as source,
        zipfile.ZipFile(tmp_path / "template-customized.pptx", "w", compression=zipfile.ZIP_DEFLATED) as target,
    ):
        for info in source.infolist():
            data = source.read(info.filename)
            if info.filename == "ppt/theme/theme1.xml":
                root = ET.fromstring(data)
                clr = root.find(".//a:clrScheme", NS)
                assert clr is not None
                clr.set("name", "Custom Template Theme")
                accent1 = clr.find("a:accent1", NS)
                assert accent1 is not None
                accent1[0].set("val", "ABCDEF")
                font_scheme = root.find(".//a:fontScheme", NS)
                assert font_scheme is not None
                major = font_scheme.find("a:majorFont/a:latin", NS)
                minor = font_scheme.find("a:minorFont/a:latin", NS)
                assert major is not None
                assert minor is not None
                major.set("typeface", "Template Headings")
                minor.set("typeface", "Template Body")
                data = ET.tostring(root, encoding="utf-8", xml_declaration=True)
            target.writestr(info, data)
    customized_template = tmp_path / "template-customized.pptx"
    deck = parse_deck(
        "# Slide\n\nBody\n",
        input_path=tmp_path / "deck.md",
        source_name=str(tmp_path / "deck.md"),
    )
    output = tmp_path / "deck.pptx"

    render_pptx(deck, output_path=output, template_path=customized_template, force=False, base_dir=tmp_path)

    with zipfile.ZipFile(output) as zf:
        theme = ET.fromstring(zf.read("ppt/theme/theme1.xml"))
    clr = theme.find(".//a:clrScheme", NS)
    assert clr is not None
    assert clr.attrib["name"] == "Custom Template Theme"
    dark1 = clr.find("a:dk1", NS)
    assert dark1 is not None
    assert dark1[0].tag == f"{{{NS['a']}}}sysClr"
    assert dark1[0].attrib == {"val": "windowText", "lastClr": "000000"}
    accent1 = clr.find("a:accent1", NS)
    assert accent1 is not None
    assert accent1[0].attrib["val"] == "ABCDEF"
    font_scheme = theme.find(".//a:fontScheme", NS)
    assert font_scheme is not None
    major = font_scheme.find("a:majorFont/a:latin", NS)
    minor = font_scheme.find("a:minorFont/a:latin", NS)
    assert major is not None
    assert minor is not None
    assert major.attrib["typeface"] == "Template Headings"
    assert minor.attrib["typeface"] == "Template Body"


def test_render_default_body_paragraph_spacing_without_template_override(tmp_path: Path) -> None:
    deck = parse_deck(
        "# Slide\n\nParagraph text.\n",
        input_path=tmp_path / "deck.md",
        source_name=str(tmp_path / "deck.md"),
    )
    output = tmp_path / "deck.pptx"

    render_pptx(deck, output_path=output, template_path=None, force=False, base_dir=tmp_path)

    with zipfile.ZipFile(output) as zf:
        slide_xml = zf.read("ppt/slides/slide1.xml").decode("utf-8")
    assert '<a:lnSpc><a:spcPct val="100000"/></a:lnSpc>' in slide_xml
    assert '<a:spcBef><a:spcPts val="1200"/></a:spcBef>' in slide_xml
    assert '<a:spcAft><a:spcPts val="600"/></a:spcAft>' in slide_xml


def test_render_preserves_template_paragraph_spacing_when_template_specified(tmp_path: Path) -> None:
    deck = parse_deck(
        "# Slide\n\nParagraph text.\n",
        input_path=tmp_path / "deck.md",
        source_name=str(tmp_path / "deck.md"),
    )
    output = tmp_path / "deck.pptx"

    render_pptx(deck, output_path=output, template_path=default_template_path(), force=False, base_dir=tmp_path)

    with zipfile.ZipFile(output) as zf:
        slide_xml = zf.read("ppt/slides/slide1.xml").decode("utf-8")
    assert "<a:lnSpc>" not in slide_xml
    assert "<a:spcBef>" not in slide_xml
    assert "<a:spcAft>" not in slide_xml


def test_render_background_and_body_image(tmp_path: Path) -> None:
    image_path = tmp_path / "photo.png"
    image_path.write_bytes(PNG_BYTES)
    deck = parse_deck(
        """---
background: "linear-gradient(90deg, #0E2841 0%, #156082 100%)"
color_scheme:
  preset: Office
---

# Photo

![alt](./photo.png)
""",
        input_path=tmp_path / "deck.md",
        source_name=str(tmp_path / "deck.md"),
    )
    output = tmp_path / "deck.pptx"

    render_pptx(deck, output_path=output, template_path=None, force=False, base_dir=tmp_path)

    prs = Presentation(str(output))
    names = [shape.name for shape in prs.slides[0].shapes]
    assert prs.slide_masters[0].background.fill.type == 3
    assert round(prs.slide_masters[0].background.fill.gradient_angle) == 90
    assert "MarkdownSlidesImage" in names


def test_render_remote_image_with_fake_downloader(tmp_path: Path) -> None:
    deck = parse_deck(
        "# Slide\n\n![alt](https://example.com/photo.png)\n",
        input_path=tmp_path / "deck.md",
        source_name=str(tmp_path / "deck.md"),
    )
    output = tmp_path / "deck.pptx"

    render_pptx(
        deck,
        output_path=output,
        template_path=None,
        force=False,
        base_dir=tmp_path,
        downloader=FakeDownloader(),
    )

    prs = Presentation(str(output))
    assert "MarkdownSlidesImage" in [shape.name for shape in prs.slides[0].shapes]


def test_render_lists_and_headings_have_expected_bullet_xml(tmp_path: Path) -> None:
    deck = parse_deck(
        "# Slide\n\n## Heading\nParagraph.\n\n- Bullet\n1. Numbered\n",
        input_path=tmp_path / "deck.md",
        source_name=str(tmp_path / "deck.md"),
    )
    output = tmp_path / "deck.pptx"

    render_pptx(deck, output_path=output, template_path=None, force=False, base_dir=tmp_path)

    with zipfile.ZipFile(output) as zf:
        slide_xml = zf.read("ppt/slides/slide1.xml").decode("utf-8")
    assert "<a:buNone/>" in slide_xml
    assert "arabicPeriod" in slide_xml
    assert "• Bullet" not in slide_xml


def test_render_document_background_image_targets_slide_master(tmp_path: Path) -> None:
    image_path = tmp_path / "bg.png"
    image_path.write_bytes(PNG_BYTES)
    deck = parse_deck(
        """---
background: "url('./bg.png')"
---

# Slide

Body
""",
        input_path=tmp_path / "deck.md",
        source_name=str(tmp_path / "deck.md"),
    )
    output = tmp_path / "deck.pptx"

    render_pptx(deck, output_path=output, template_path=None, force=False, base_dir=tmp_path)

    with zipfile.ZipFile(output) as zf:
        master_xml = zf.read("ppt/slideMasters/slideMaster1.xml").decode("utf-8")
        slide_xml = zf.read("ppt/slides/slide1.xml").decode("utf-8")
    assert "<a:blipFill>" in master_xml
    assert 'r:embed="' in master_xml
    assert "MarkdownSlidesBackgroundImage" not in slide_xml


def test_render_theme_font_refs_blockquote_indent_and_table_style(tmp_path: Path) -> None:
    deck = parse_deck(
        """# Formatting

## Heading
Paragraph

> Quote text

# Table

| A | B |
| --- | --- |
| 1 | 2 |
""",
        input_path=tmp_path / "deck.md",
        source_name=str(tmp_path / "deck.md"),
    )
    output = tmp_path / "deck.pptx"

    render_pptx(deck, output_path=output, template_path=None, force=False, base_dir=tmp_path)

    with zipfile.ZipFile(output) as zf:
        slide_xml = zf.read("ppt/slides/slide1.xml").decode("utf-8")
        table_slide_xml = zf.read("ppt/slides/slide2.xml").decode("utf-8")
    assert "+mj-lt" in slide_xml
    assert "+mn-lt" in slide_xml
    assert 'marL="0"' in slide_xml
    assert 'indent="0"' in slide_xml
    assert 'schemeClr val="accent1"' in slide_xml
    assert TABLE_STYLE_MEDIUM_1_ACCENT_1 in table_slide_xml
    table_properties = ET.fromstring(table_slide_xml).find(".//a:tblPr", NS)
    assert table_properties is not None
    assert table_properties.attrib == {"firstRow": "1", "bandRow": "1"}


def test_render_slide_table_options_as_native_powerpoint_flags(tmp_path: Path) -> None:
    deck = parse_deck(
        """# Table
---
table:
  header_row: false
  total_row: true
  first_column: true
  last_column: true
  banded_rows: false
  banded_columns: true
---

| A | B |
| --- | --- |
| 1 | 2 |
""",
        input_path=tmp_path / "deck.md",
        source_name=str(tmp_path / "deck.md"),
    )
    output = tmp_path / "deck.pptx"

    render_pptx(deck, output_path=output, template_path=None, force=False, base_dir=tmp_path)

    with zipfile.ZipFile(output) as zf:
        slide_xml = zf.read("ppt/slides/slide1.xml").decode("utf-8")
    table_properties = ET.fromstring(slide_xml).find(".//a:tblPr", NS)
    assert table_properties is not None
    assert table_properties.attrib == {
        "lastRow": "1",
        "firstCol": "1",
        "lastCol": "1",
        "bandCol": "1",
    }


def test_render_radial_gradient_background_writes_path_gradient_xml(tmp_path: Path) -> None:
    deck = parse_deck(
        """# Radial
---
background: "radial-gradient(circle, #0E2841 0%, #156082 55%, #EAF3FF 100%)"
---

Body
""",
        input_path=tmp_path / "deck.md",
        source_name=str(tmp_path / "deck.md"),
    )
    output = tmp_path / "deck.pptx"

    render_pptx(deck, output_path=output, template_path=None, force=False, base_dir=tmp_path)

    with zipfile.ZipFile(output) as zf:
        slide_xml = zf.read("ppt/slides/slide1.xml").decode("utf-8")
    assert '<a:path path="circle">' in slide_xml
    assert '<a:fillToRect l="50000" t="50000" r="50000" b="50000"/>' in slide_xml


def test_render_theme_color_refs_for_text_and_backgrounds(tmp_path: Path) -> None:
    deck = parse_deck(
        """---
title_color: "var(--light-1)"
body_color: "var(--dark-1)"
background: "linear-gradient(90deg, var(--accent-1) 0%, var(--accent-2) 100%)"
---

# Slide

Paragraph

> Quote text
""",
        input_path=tmp_path / "deck.md",
        source_name=str(tmp_path / "deck.md"),
    )
    output = tmp_path / "deck.pptx"

    render_pptx(deck, output_path=output, template_path=None, force=False, base_dir=tmp_path)

    with zipfile.ZipFile(output) as zf:
        master_xml = zf.read("ppt/slideMasters/slideMaster1.xml").decode("utf-8")
        slide_xml = zf.read("ppt/slides/slide1.xml").decode("utf-8")
    assert 'schemeClr val="accent1"' in master_xml
    assert 'schemeClr val="accent2"' in master_xml
    assert 'schemeClr val="lt1"' in slide_xml
    assert 'schemeClr val="dk1"' in slide_xml
    assert 'schemeClr val="accent1"' in slide_xml


def test_render_section_header_subtitle_uses_theme_body_text_color_by_default(tmp_path: Path) -> None:
    deck = parse_deck(
        """# Section heading
---
layout: Section Header
---

Subtitle text
""",
        input_path=tmp_path / "deck.md",
        source_name=str(tmp_path / "deck.md"),
    )
    output = tmp_path / "deck.pptx"

    render_pptx(deck, output_path=output, template_path=None, force=False, base_dir=tmp_path)

    with zipfile.ZipFile(output) as zf:
        slide_xml = zf.read("ppt/slides/slide1.xml").decode("utf-8")
    assert 'schemeClr val="dk1"' in slide_xml
    assert "<a:defRPr" in slide_xml
    assert "<a:endParaRPr" in slide_xml


def test_force_flag_controls_overwrite(tmp_path: Path) -> None:
    deck = parse_deck("# Slide\n\nBody\n", input_path=tmp_path / "deck.md", source_name=str(tmp_path / "deck.md"))
    output = tmp_path / "deck.pptx"
    render_pptx(deck, output_path=output, template_path=None, force=False, base_dir=tmp_path)

    with pytest.raises(RenderError):
        render_pptx(deck, output_path=output, template_path=None, force=False, base_dir=tmp_path)

    render_pptx(deck, output_path=output, template_path=None, force=True, base_dir=tmp_path)


def test_output_path_must_be_a_file_even_with_force(tmp_path: Path) -> None:
    deck = parse_deck("# Slide\n\nBody\n", input_path=tmp_path / "deck.md", source_name=str(tmp_path / "deck.md"))
    output = tmp_path / "existing-directory"
    output.mkdir()

    with pytest.raises(RenderError) as excinfo:
        render_pptx(deck, output_path=output, template_path=None, force=True, base_dir=tmp_path)

    assert excinfo.value.context.code == "output_not_file"
    assert list(output.iterdir()) == []
    assert output.exists()


def test_missing_placeholder_template_fails(tmp_path: Path) -> None:
    template = tmp_path / "broken-template.pptx"
    with (
        zipfile.ZipFile(default_template_path(), "r") as source,
        zipfile.ZipFile(template, "w", compression=zipfile.ZIP_DEFLATED) as target,
    ):
        for info in source.infolist():
            data = source.read(info.filename)
            if info.filename == "ppt/slideLayouts/slideLayout2.xml":
                xml = ET.fromstring(data)
                sp_tree = xml.find(".//p:spTree", {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"})
                for sp in list(sp_tree):
                    ph = sp.find(".//p:ph", {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"})
                    if ph is not None and ph.attrib.get("idx") == "1":
                        sp_tree.remove(sp)
                        break
                data = ET.tostring(xml, encoding="utf-8", xml_declaration=True)
            target.writestr(info, data)
    deck = parse_deck("# Slide\n\nBody\n", input_path=tmp_path / "deck.md", source_name=str(tmp_path / "deck.md"))
    output = tmp_path / "deck.pptx"

    with pytest.raises(TemplateError):
        render_pptx(deck, output_path=output, template_path=template, force=False, base_dir=tmp_path)


def test_render_preserves_all_gradient_stops_and_zero_degree_angle(tmp_path: Path) -> None:
    deck = parse_deck(
        """---
background: "linear-gradient(0deg, #111111 0%, #222222 40%, #333333 100%)"
---

# Slide
---
background: "linear-gradient(0deg, #444444 0%, #555555 50%, #666666 100%)"
---

Body
""",
        input_path=tmp_path / "deck.md",
        source_name=str(tmp_path / "deck.md"),
    )
    output = tmp_path / "deck.pptx"

    render_pptx(deck, output_path=output, template_path=None, force=False, base_dir=tmp_path)

    with zipfile.ZipFile(output) as zf:
        master = ET.fromstring(zf.read("ppt/slideMasters/slideMaster1.xml"))
        slide = ET.fromstring(zf.read("ppt/slides/slide1.xml"))
    master_stops = master.findall(
        ".//p:bg/p:bgPr/a:gradFill/a:gsLst/a:gs",
        {**NS, "p": "http://schemas.openxmlformats.org/presentationml/2006/main"},
    )
    slide_stops = slide.findall(
        ".//p:bg/p:bgPr/a:gradFill/a:gsLst/a:gs",
        {**NS, "p": "http://schemas.openxmlformats.org/presentationml/2006/main"},
    )
    assert [stop.attrib["pos"] for stop in master_stops] == ["0", "40000", "100000"]
    assert [stop.attrib["pos"] for stop in slide_stops] == ["0", "50000", "100000"]
    assert (
        master.find(
            ".//p:bg/p:bgPr/a:gradFill/a:lin", {**NS, "p": "http://schemas.openxmlformats.org/presentationml/2006/main"}
        ).attrib["ang"]
        == "0"
    )
    assert (
        slide.find(
            ".//p:bg/p:bgPr/a:gradFill/a:lin", {**NS, "p": "http://schemas.openxmlformats.org/presentationml/2006/main"}
        ).attrib["ang"]
        == "0"
    )


def test_layout_details_report_unsupported_layouts() -> None:
    details = list_layout_details()

    assert details["master"]["index"] == 1
    title_content = next(item for item in details["layouts"] if item["name"] == "Title and Content")
    assert title_content["compatible"] is True
    two_content = next(item for item in details["layouts"] if item["name"] == "Two Content")
    assert two_content["compatible"] is False


def test_ambiguous_placeholder_template_fails(tmp_path: Path) -> None:
    template = tmp_path / "ambiguous-template.pptx"
    with (
        zipfile.ZipFile(default_template_path(), "r") as source,
        zipfile.ZipFile(template, "w", compression=zipfile.ZIP_DEFLATED) as target,
    ):
        for info in source.infolist():
            data = source.read(info.filename)
            if info.filename == "ppt/slideLayouts/slideLayout2.xml":
                xml = ET.fromstring(data)
                sp_tree = xml.find(".//p:spTree", {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"})
                body_shape = next(
                    shape
                    for shape in list(sp_tree)
                    if (
                        ph := shape.find(".//p:ph", {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"})
                    )
                    is not None
                    and ph.attrib.get("idx") == "1"
                )
                duplicate = deepcopy(body_shape)
                c_nv_pr = duplicate.find(
                    ".//p:cNvPr", {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}
                )
                c_nv_pr.set("id", "999")
                c_nv_pr.set("name", "Duplicate Content Placeholder")
                sp_tree.append(duplicate)
                data = ET.tostring(xml, encoding="utf-8", xml_declaration=True)
            target.writestr(info, data)
    deck = parse_deck("# Slide\n\nBody\n", input_path=tmp_path / "deck.md", source_name=str(tmp_path / "deck.md"))

    with pytest.raises(TemplateError) as excinfo:
        render_pptx(
            deck,
            output_path=tmp_path / "deck.pptx",
            template_path=template,
            force=False,
            base_dir=tmp_path,
        )

    assert excinfo.value.context.code == "ambiguous_placeholder"


def test_duplicate_layout_name_template_fails(tmp_path: Path) -> None:
    template = tmp_path / "duplicate-layout-template.pptx"
    with (
        zipfile.ZipFile(default_template_path(), "r") as source,
        zipfile.ZipFile(template, "w", compression=zipfile.ZIP_DEFLATED) as target,
    ):
        for info in source.infolist():
            data = source.read(info.filename)
            if info.filename == "ppt/slideLayouts/slideLayout6.xml":
                xml = ET.fromstring(data)
                c_sld = xml.find(".//p:cSld", {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"})
                c_sld.set("name", "Title and Content")
                data = ET.tostring(xml, encoding="utf-8", xml_declaration=True)
            target.writestr(info, data)
    deck = parse_deck("# Slide\n\nBody\n", input_path=tmp_path / "deck.md", source_name=str(tmp_path / "deck.md"))

    with pytest.raises(TemplateError) as excinfo:
        render_pptx(
            deck,
            output_path=tmp_path / "deck.pptx",
            template_path=template,
            force=False,
            base_dir=tmp_path,
        )

    assert excinfo.value.context.code == "ambiguous_layout"


def test_downloader_streams_validates_and_caches_remote_images() -> None:
    requests = 0

    def handler(request: httpx.Request) -> httpx.Response:
        nonlocal requests
        requests += 1
        return httpx.Response(200, headers={"content-type": "image/png"}, content=PNG_BYTES)

    client = httpx.Client(transport=httpx.MockTransport(handler))
    downloader = Downloader(client=client)
    try:
        assert downloader.fetch("https://example.com/image.png") == PNG_BYTES
        assert downloader.fetch("https://example.com/image.png") == PNG_BYTES
    finally:
        client.close()
    assert requests == 1


def test_downloader_rejects_non_image_content_type() -> None:
    client = httpx.Client(
        transport=httpx.MockTransport(
            lambda request: httpx.Response(200, headers={"content-type": "text/html"}, content=b"not an image")
        )
    )
    downloader = Downloader(client=client)
    try:
        with pytest.raises(AssetError) as excinfo:
            downloader.fetch("https://user:secret@example.com/image.png?token=secret")
    finally:
        client.close()
    assert excinfo.value.context.code == "invalid_remote_image_type"
    assert "secret" not in excinfo.value.context.message


def test_downloader_rejects_declared_size_over_limit() -> None:
    client = httpx.Client(transport=httpx.MockTransport(lambda request: httpx.Response(200, content=PNG_BYTES)))
    downloader = Downloader(client=client, max_bytes=8)
    try:
        with pytest.raises(AssetError) as excinfo:
            downloader.fetch("https://example.com/image.png")
    finally:
        client.close()
    assert excinfo.value.context.code == "remote_image_too_large"


def test_downloader_enforces_streaming_size_limit_without_content_length() -> None:
    client = httpx.Client(
        transport=httpx.MockTransport(
            lambda request: httpx.Response(
                200,
                headers={"content-type": "image/png"},
                stream=httpx.ByteStream(PNG_BYTES),
            )
        )
    )
    downloader = Downloader(client=client, max_bytes=8)
    try:
        with pytest.raises(AssetError) as excinfo:
            downloader.fetch("https://example.com/image.png")
    finally:
        client.close()
    assert excinfo.value.context.code == "remote_image_too_large"


def test_downloader_sanitizes_http_errors() -> None:
    client = httpx.Client(transport=httpx.MockTransport(lambda request: httpx.Response(404)))
    downloader = Downloader(client=client)
    try:
        with pytest.raises(AssetError) as excinfo:
            downloader.fetch("https://user:password@example.com/private/token.png?api_key=secret")
    finally:
        client.close()
    assert excinfo.value.context.code == "image_download_failed"
    assert "404" in excinfo.value.context.message
    assert "password" not in excinfo.value.context.message
    assert "token" not in excinfo.value.context.message
    assert "secret" not in excinfo.value.context.message


def test_render_can_disable_remote_images(tmp_path: Path) -> None:
    deck = parse_deck(
        "# Slide\n\n![alt](https://example.com/photo.png)\n",
        input_path=tmp_path / "deck.md",
        source_name=str(tmp_path / "deck.md"),
    )

    downloader = FakeDownloader()
    with pytest.raises(AssetError) as excinfo:
        render_pptx(
            deck,
            output_path=tmp_path / "deck.pptx",
            template_path=None,
            force=False,
            base_dir=tmp_path,
            downloader=downloader,
            allow_remote_images=False,
        )

    assert excinfo.value.context.code == "remote_images_disabled"
    assert downloader.enabled is True
    downloader.close()


def test_render_rejects_remote_bytes_that_are_not_an_image(tmp_path: Path) -> None:
    class InvalidImageDownloader(FakeDownloader):
        def fetch(self, url: str) -> bytes:
            return b"not an image"

    deck = parse_deck(
        "# Slide\n\n![alt](https://example.com/photo.png)\n",
        input_path=tmp_path / "deck.md",
        source_name=str(tmp_path / "deck.md"),
    )
    downloader = InvalidImageDownloader()
    try:
        with pytest.raises(AssetError) as excinfo:
            render_pptx(
                deck,
                output_path=tmp_path / "deck.pptx",
                template_path=None,
                force=False,
                base_dir=tmp_path,
                downloader=downloader,
            )
    finally:
        downloader.close()

    assert excinfo.value.context.code == "invalid_image"


def test_render_rejects_images_over_the_pixel_limit(tmp_path: Path) -> None:
    image_path = tmp_path / "oversized.png"
    image_path.write_bytes(_png_header(width=10_000, height=6_000))
    deck = parse_deck(
        "# Slide\n\n![alt](./oversized.png)\n",
        input_path=tmp_path / "deck.md",
        source_name=str(tmp_path / "deck.md"),
    )

    with pytest.raises(AssetError) as excinfo:
        render_pptx(
            deck,
            output_path=tmp_path / "deck.pptx",
            template_path=None,
            force=False,
            base_dir=tmp_path,
        )

    assert excinfo.value.context.code == "image_dimensions_too_large"


def test_render_fully_custom_color_scheme_metadata(tmp_path: Path) -> None:
    deck = parse_deck(
        """---
color_scheme:
  preset:
  dark_1: "#010101"
  light_1: "#F1F1F1"
  dark_2: "#020202"
  light_2: "#F2F2F2"
  accent_1: "#111111"
  accent_2: "#222222"
  accent_3: "#333333"
  accent_4: "#444444"
  accent_5: "#555555"
  accent_6: "#666666"
  hyperlink: "#777777"
  followed_hyperlink: "#888888"
---

# Slide

Body
""",
        input_path=tmp_path / "deck.md",
        source_name=str(tmp_path / "deck.md"),
    )
    output = tmp_path / "deck.pptx"

    render_pptx(deck, output_path=output, template_path=None, force=False, base_dir=tmp_path)

    with zipfile.ZipFile(output) as zf:
        theme = ET.fromstring(zf.read("ppt/theme/theme1.xml"))
    color_scheme = theme.find(".//a:clrScheme", NS)
    assert color_scheme is not None
    assert color_scheme.attrib["name"] == "Custom"
    expected_colors = {
        "dk1": "010101",
        "lt1": "F1F1F1",
        "dk2": "020202",
        "lt2": "F2F2F2",
        "accent1": "111111",
        "accent2": "222222",
        "accent3": "333333",
        "accent4": "444444",
        "accent5": "555555",
        "accent6": "666666",
        "hlink": "777777",
        "folHlink": "888888",
    }
    for xml_key, expected_color in expected_colors.items():
        parent = color_scheme.find(f"a:{xml_key}", NS)
        assert parent is not None
        assert len(parent) == 1
        assert parent[0].tag == f"{{{NS['a']}}}srgbClr"
        assert parent[0].attrib == {"val": expected_color}


def test_render_rewrites_the_first_masters_actual_theme_part(tmp_path: Path) -> None:
    template = tmp_path / "renamed-theme-template.pptx"
    relationship_ns = "http://schemas.openxmlformats.org/package/2006/relationships"
    with (
        zipfile.ZipFile(default_template_path(), "r") as source,
        zipfile.ZipFile(template, "w", compression=zipfile.ZIP_DEFLATED) as target,
    ):
        for info in source.infolist():
            data = source.read(info.filename)
            if info.filename == "ppt/theme/theme1.xml":
                target.writestr("ppt/theme/theme7.xml", data)
                continue
            if info.filename == "ppt/slideMasters/_rels/slideMaster1.xml.rels":
                xml = ET.fromstring(data)
                theme_rel = next(
                    rel
                    for rel in xml.findall(f"{{{relationship_ns}}}Relationship")
                    if rel.attrib["Type"].endswith("/theme")
                )
                theme_rel.set("Target", "../theme/theme7.xml")
                data = ET.tostring(xml, encoding="utf-8", xml_declaration=True)
            target.writestr(info, data)
    deck = parse_deck(
        """---
color_scheme:
  preset: Blue Warm
---

# Slide

Body
""",
        input_path=tmp_path / "deck.md",
        source_name=str(tmp_path / "deck.md"),
    )
    output = tmp_path / "deck.pptx"

    render_pptx(deck, output_path=output, template_path=template, force=False, base_dir=tmp_path)

    with zipfile.ZipFile(output) as zf:
        assert "ppt/theme/theme1.xml" not in zf.namelist()
        theme = ET.fromstring(zf.read("ppt/theme/theme7.xml"))
    color_scheme = theme.find(".//a:clrScheme", NS)
    assert color_scheme.attrib["name"] == "Blue Warm"


def test_render_retains_all_template_masters(tmp_path: Path) -> None:
    template = tmp_path / "two-master-template.pptx"
    _make_two_master_template(template)
    assert len(Presentation(str(template)).slide_masters) == 2
    deck = parse_deck("# Slide\n\nBody\n", input_path=tmp_path / "deck.md", source_name=str(tmp_path / "deck.md"))
    output = tmp_path / "deck.pptx"

    render_pptx(deck, output_path=output, template_path=template, force=False, base_dir=tmp_path)

    rendered = Presentation(str(output))
    assert len(rendered.slide_masters) == 2
    assert rendered.slides[0].slide_layout.part.partname == "/ppt/slideLayouts/slideLayout2.xml"
    with zipfile.ZipFile(output) as zf:
        assert "ppt/slideMasters/slideMaster2.xml" in zf.namelist()


def test_list_and_select_template_master(tmp_path: Path) -> None:
    template = tmp_path / "two-master-template.pptx"
    _make_two_master_template(template)

    masters = list_master_details(template)
    layouts = list_layout_details(template, master="Executive Theme")

    assert [master["index"] for master in masters] == [1, 2]
    assert masters[1]["name"] == "Executive"
    assert masters[1]["theme_name"] == "Executive Theme"
    assert "Executive" in masters[1]["selectable_names"]
    assert layouts["master"]["index"] == 2
    assert any(layout["name"] == "Title and Content" for layout in layouts["layouts"])


def test_cli_master_selection_and_slide_override_choose_layouts_from_each_master(tmp_path: Path) -> None:
    template = tmp_path / "two-master-template.pptx"
    _make_two_master_template(template)
    deck = parse_deck(
        """# Default master

Body

# Slide override
---
master: Executive
---

Body
""",
        input_path=tmp_path / "deck.md",
        source_name=str(tmp_path / "deck.md"),
    )
    output = tmp_path / "deck.pptx"
    report: dict[str, object] = {}

    render_pptx(
        deck,
        output_path=output,
        template_path=template,
        force=False,
        base_dir=tmp_path,
        master=1,
        report=report,
    )

    rendered = Presentation(str(output))
    assert rendered.slides[0].slide_layout.part.partname == "/ppt/slideLayouts/slideLayout2.xml"
    assert rendered.slides[1].slide_layout.part.partname == "/ppt/slideLayouts/slideLayout12.xml"
    assert rendered.slides[1].slide_layout.slide_master.name == "Executive"
    assert report["retained_master_count"] == 2
    assert [item["index"] for item in report["masters_used"]] == [1, 2]


def test_default_master_can_be_selected_by_theme_name(tmp_path: Path) -> None:
    template = tmp_path / "two-master-template.pptx"
    _make_two_master_template(template)
    deck = parse_deck("# Slide\n\nBody\n", input_path=tmp_path / "deck.md", source_name=str(tmp_path / "deck.md"))
    output = tmp_path / "deck.pptx"

    render_pptx(
        deck,
        output_path=output,
        template_path=template,
        force=False,
        base_dir=tmp_path,
        master="Executive Theme",
    )

    rendered = Presentation(str(output))
    assert rendered.slides[0].slide_layout.part.partname == "/ppt/slideLayouts/slideLayout12.xml"


def test_master_selection_reports_invalid_index(tmp_path: Path) -> None:
    template = tmp_path / "two-master-template.pptx"
    _make_two_master_template(template)
    deck = parse_deck("# Slide\n\nBody\n", input_path=tmp_path / "deck.md", source_name=str(tmp_path / "deck.md"))

    with pytest.raises(TemplateError) as excinfo:
        render_pptx(
            deck,
            output_path=tmp_path / "deck.pptx",
            template_path=template,
            force=False,
            base_dir=tmp_path,
            master=3,
        )

    assert excinfo.value.context.code == "master_not_found"


def test_document_theme_and_background_overrides_apply_to_all_retained_masters(tmp_path: Path) -> None:
    template = tmp_path / "two-master-template.pptx"
    _make_two_master_template(template)
    deck = parse_deck(
        """---
fonts:
  body: Arial
  headings: Arial
color_scheme:
  preset: Blue Warm
background: "#EAF3FF"
---

# Slide

Body
""",
        input_path=tmp_path / "deck.md",
        source_name=str(tmp_path / "deck.md"),
    )
    output = tmp_path / "deck.pptx"

    render_pptx(deck, output_path=output, template_path=template, force=False, base_dir=tmp_path)

    with zipfile.ZipFile(output) as zf:
        for theme_name in ("ppt/theme/theme1.xml", "ppt/theme/theme7.xml"):
            theme = ET.fromstring(zf.read(theme_name))
            color_scheme = theme.find(".//a:clrScheme", NS)
            assert color_scheme.attrib["name"] == "Blue Warm"
            assert color_scheme.find("a:dk1", NS)[0].tag == f"{{{NS['a']}}}srgbClr"
            assert color_scheme.find("a:dk1", NS)[0].attrib == {"val": "000000"}
            assert color_scheme.find("a:lt1", NS)[0].tag == f"{{{NS['a']}}}srgbClr"
            assert color_scheme.find("a:lt1", NS)[0].attrib == {"val": "FFFFFF"}
            assert theme.find(".//a:majorFont/a:latin", NS).attrib["typeface"] == "Arial"
            assert theme.find(".//a:minorFont/a:latin", NS).attrib["typeface"] == "Arial"
        for master_name in ("ppt/slideMasters/slideMaster1.xml", "ppt/slideMasters/slideMaster2.xml"):
            master = ET.fromstring(zf.read(master_name))
            assert (
                master.find(".//p:bg", {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}) is not None
            )
